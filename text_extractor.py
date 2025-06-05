"""
Advanced Text Extraction Engine
==============================

A comprehensive text extraction system supporting multiple OCR engines,
document formats, image processing, and AI-powered text analysis.

Features:
- Multi-engine OCR (Tesseract, EasyOCR, PaddleOCR)
- Advanced document processing (PDF, Word, Excel, PowerPoint)
- Intelligent image preprocessing
- Language detection and translation
- Confidence scoring and quality assessment
- Batch processing with progress tracking
- Extraction history and metadata storage
"""

import os
import sys
import json
import sqlite3
import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Tuple, Union, Any
from dataclasses import dataclass
from concurrent.futures import ThreadPoolExecutor, as_completed
import tempfile
import hashlib

# Core libraries
import cv2
import numpy as np
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract

# Optional advanced OCR engines
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except (ImportError, AttributeError, Exception) as e:
    EASYOCR_AVAILABLE = False
    print(f"Warning: EasyOCR not available: {e}")

try:
    from paddleocr import PaddleOCR
    PADDLEOCR_AVAILABLE = True
except (ImportError, AttributeError, Exception) as e:
    PADDLEOCR_AVAILABLE = False
    print(f"Warning: PaddleOCR not available: {e}")

# Document processing
import PyPDF2
import pdfplumber
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except (ImportError, AttributeError, Exception) as e:
    PYMUPDF_AVAILABLE = False
    print(f"Warning: PyMuPDF not available: {e}")

from docx import Document
import openpyxl
from pptx import Presentation

# Optional xlrd for old Excel files
try:
    import xlrd
    XLRD_AVAILABLE = True
except (ImportError, AttributeError, Exception) as e:
    XLRD_AVAILABLE = False
    print(f"Warning: xlrd not available: {e}")

# Web and text processing
import requests
from bs4 import BeautifulSoup
try:
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options as ChromeOptions
    from selenium.webdriver.common.by import By
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    SELENIUM_AVAILABLE = True
except (ImportError, AttributeError, Exception) as e:
    SELENIUM_AVAILABLE = False
    print(f"Warning: Selenium not available: {e}")

# Language processing
try:
    import langdetect
    LANGDETECT_AVAILABLE = True
except (ImportError, AttributeError, Exception) as e:
    LANGDETECT_AVAILABLE = False
    print(f"Warning: langdetect not available: {e}")

try:
    from googletrans import Translator
    GOOGLETRANS_AVAILABLE = True
except (ImportError, AttributeError, Exception) as e:
    GOOGLETRANS_AVAILABLE = False
    print(f"Warning: googletrans not available: {e}")

# Image processing (optional for advanced features)
try:
    from skimage import filters, morphology, measure
    SKIMAGE_AVAILABLE = True
except ImportError:
    SKIMAGE_AVAILABLE = False

try:
    from scipy import ndimage
    SCIPY_AVAILABLE = True
except ImportError:
    SCIPY_AVAILABLE = False

# Utilities
try:
    from loguru import logger
    LOGURU_AVAILABLE = True
except ImportError:
    LOGURU_AVAILABLE = False    # Fallback to standard logging
    import logging
    logger = logging.getLogger(__name__)
    logging.basicConfig(level=logging.INFO)

try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False

try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False


@dataclass
class ExtractionResult:
    """Data class for extraction results"""
    text: str
    confidence: float
    language: str
    metadata: Dict[str, Any]
    processing_time: float
    source_file: str
    extraction_method: str
    timestamp: datetime


class ImagePreprocessor:
    """Advanced image preprocessing for better OCR results"""
    
    def __init__(self):
        self.techniques = {
            'denoise': self._denoise,
            'deskew': self._deskew,
            'enhance_contrast': self._enhance_contrast,
            'binarize': self._binarize,
            'remove_shadows': self._remove_shadows,
            'sharpen': self._sharpen,
            'resize': self._resize_for_ocr
        }
    
    def preprocess(self, image: np.ndarray, techniques: List[str] = None) -> np.ndarray:
        """Apply preprocessing techniques to improve OCR accuracy"""
        if techniques is None:
            techniques = ['denoise', 'deskew', 'enhance_contrast', 'binarize']
        
        processed = image.copy()
        
        for technique in techniques:
            if technique in self.techniques:
                try:
                    processed = self.techniques[technique](processed)
                    logger.debug(f"Applied {technique} preprocessing")
                except Exception as e:
                    logger.warning(f"Failed to apply {technique}: {e}")
                    
        return processed
    
    def _denoise(self, image: np.ndarray) -> np.ndarray:
        """Remove noise from image"""
        return cv2.fastNlMeansDenoising(image)
    
    def _deskew(self, image: np.ndarray) -> np.ndarray:
        """Correct skew in the image"""
        # Convert to grayscale if needed
        if len(image.shape) == 3:
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        else:
            gray = image.copy()
        
        # Find edges
        edges = cv2.Canny(gray, 50, 150, apertureSize=3)
        
        # Detect lines
        lines = cv2.HoughLines(edges, 1, np.pi/180, threshold=100)
        
        if lines is not None:
            # Calculate average angle
            angles = []
            for line in lines:
                rho, theta = line[0]
                angle = theta * 180 / np.pi
                if angle < 45:
                    angles.append(angle)
                elif angle > 135:
                    angles.append(angle - 180)
            
            if angles:
                avg_angle = np.mean(angles)
                
                # Rotate image
                h, w = image.shape[:2]
                center = (w // 2, h // 2)
                rotation_matrix = cv2.getRotationMatrix2D(center, avg_angle, 1.0)
                rotated = cv2.warpAffine(image, rotation_matrix, (w, h), 
                                       flags=cv2.INTER_CUBIC, 
                                       borderMode=cv2.BORDER_REPLICATE)
                return rotated
        
        return image
    
    def _enhance_contrast(self, image: np.ndarray) -> np.ndarray:
        """Enhance image contrast using CLAHE"""
        if len(image.shape) == 3:
            # Convert to LAB color space
            lab = cv2.cvtColor(image, cv2.COLOR_BGR2LAB)
            l, a, b = cv2.split(lab)
            
            # Apply CLAHE to L channel
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            l = clahe.apply(l)
            
            # Merge channels and convert back
            enhanced = cv2.merge([l, a, b])
            return cv2.cvtColor(enhanced, cv2.COLOR_LAB2BGR)
        else:
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            return clahe.apply(image)
    
    def _binarize(self, image: np.ndarray) -> np.ndarray:
        """Convert to binary image using adaptive thresholding"""
        if len(image.shape) == 3:
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        else:
            gray = image.copy()
        
        # Apply Gaussian blur
        blurred = cv2.GaussianBlur(gray, (5, 5), 0)
        
        # Adaptive thresholding
        binary = cv2.adaptiveThreshold(
            blurred, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
            cv2.THRESH_BINARY, 11, 2
        )
        
        return binary
    
    def _remove_shadows(self, image: np.ndarray) -> np.ndarray:
        """Remove shadows from image"""
        if len(image.shape) == 3:
            # Convert to LAB
            lab = cv2.cvtColor(image, cv2.COLOR_BGR2LAB)
            l, a, b = cv2.split(lab)
            
            # Apply morphological operations to L channel
            kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (20, 20))
            opened = cv2.morphologyEx(l, cv2.MORPH_OPEN, kernel)
            normalized = cv2.divide(l, opened, scale=255)
            
            # Merge and convert back
            result = cv2.merge([normalized, a, b])
            return cv2.cvtColor(result, cv2.COLOR_LAB2BGR)
        
        return image
    
    def _sharpen(self, image: np.ndarray) -> np.ndarray:
        """Sharpen image for better text clarity"""
        kernel = np.array([[-1, -1, -1],
                          [-1,  9, -1],
                          [-1, -1, -1]])
        return cv2.filter2D(image, -1, kernel)
    
    def _resize_for_ocr(self, image: np.ndarray) -> np.ndarray:
        """Resize image to optimal size for OCR"""
        h, w = image.shape[:2]
        
        # Target height for good OCR performance
        target_height = 800
        
        if h < target_height:
            # Upscale small images
            scale = target_height / h
            new_w = int(w * scale)
            new_h = int(h * scale)
            
            return cv2.resize(image, (new_w, new_h), interpolation=cv2.INTER_CUBIC)
        elif h > 2000:
            # Downscale very large images
            scale = 2000 / h
            new_w = int(w * scale)
            new_h = int(h * scale)
            
            return cv2.resize(image, (new_w, new_h), interpolation=cv2.INTER_AREA)
        
        return image


class MultiEngineOCR:
    """Multi-engine OCR system with automatic engine selection"""
    
    def __init__(self, config: Dict[str, Any] = None):
        self.config = config or {}
        self.preprocessor = ImagePreprocessor()
        
        # Initialize available engines
        self.engines = {}
        self._init_tesseract()
        self._init_easyocr()
        self._init_paddleocr()
        
        logger.info(f"Initialized OCR engines: {list(self.engines.keys())}")
    
    def _init_tesseract(self):
        """Initialize Tesseract OCR"""
        try:
            # Test Tesseract availability
            pytesseract.get_tesseract_version()
            self.engines['tesseract'] = self._ocr_tesseract
            logger.info("Tesseract OCR initialized")
        except Exception as e:
            logger.warning(f"Tesseract not available: {e}")
    
    def _init_easyocr(self):
        """Initialize EasyOCR"""
        if EASYOCR_AVAILABLE:
            try:
                self.easyocr_reader = easyocr.Reader(['en'])
                self.engines['easyocr'] = self._ocr_easyocr
                logger.info("EasyOCR initialized")
            except Exception as e:
                logger.warning(f"EasyOCR initialization failed: {e}")
    
    def _init_paddleocr(self):
        """Initialize PaddleOCR"""
        if PADDLEOCR_AVAILABLE:
            try:
                self.paddleocr_reader = PaddleOCR(use_angle_cls=True, lang='en')
                self.engines['paddleocr'] = self._ocr_paddleocr
                logger.info("PaddleOCR initialized")
            except Exception as e:
                logger.warning(f"PaddleOCR initialization failed: {e}")
    
    def extract_text(self, image: np.ndarray, engine: str = 'auto', 
                    preprocess: bool = True) -> ExtractionResult:
        """Extract text using specified or best available engine"""
        start_time = datetime.now()
        
        # Preprocess image if requested
        if preprocess:
            processed_image = self.preprocessor.preprocess(image)
        else:
            processed_image = image
        
        # Select engine
        if engine == 'auto':
            engine = self._select_best_engine(processed_image)
        
        if engine not in self.engines:
            raise ValueError(f"Engine '{engine}' not available. Available: {list(self.engines.keys())}")
        
        # Perform OCR
        try:
            text, confidence = self.engines[engine](processed_image)
            
            # Detect language
            language = self._detect_language(text)
            
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return ExtractionResult(
                text=text,
                confidence=confidence,
                language=language,
                metadata={
                    'engine': engine,
                    'preprocessed': preprocess,
                    'image_shape': image.shape
                },
                processing_time=processing_time,
                source_file='',
                extraction_method='ocr',
                timestamp=datetime.now()
            )
            
        except Exception as e:
            logger.error(f"OCR extraction failed with {engine}: {e}")
            raise
    
    def _select_best_engine(self, image: np.ndarray) -> str:
        """Select the best OCR engine based on image characteristics"""
        # Simple heuristics for engine selection
        h, w = image.shape[:2]
        
        # For small images, prefer EasyOCR
        if h < 100 or w < 100:
            if 'easyocr' in self.engines:
                return 'easyocr'
        
        # For large images with potentially complex layouts, prefer PaddleOCR
        if h > 1000 and w > 1000:
            if 'paddleocr' in self.engines:
                return 'paddleocr'
        
        # Default to Tesseract
        if 'tesseract' in self.engines:
            return 'tesseract'
        
        # Return any available engine
        return list(self.engines.keys())[0]
    
    def _ocr_tesseract(self, image: np.ndarray) -> Tuple[str, float]:
        """Perform OCR using Tesseract"""
        # Get text and confidence data
        data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
        
        # Calculate confidence
        confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
        avg_confidence = sum(confidences) / len(confidences) if confidences else 0
        
        # Extract text
        text = pytesseract.image_to_string(image).strip()
        
        return text, avg_confidence / 100.0
    
    def _ocr_easyocr(self, image: np.ndarray) -> Tuple[str, float]:
        """Perform OCR using EasyOCR"""
        results = self.easyocr_reader.readtext(image)
        
        if not results:
            return "", 0.0
        
        # Extract text and calculate average confidence
        texts = []
        confidences = []
        
        for (bbox, text, conf) in results:
            texts.append(text)
            confidences.append(conf)
        
        combined_text = ' '.join(texts)
        avg_confidence = sum(confidences) / len(confidences) if confidences else 0
        
        return combined_text, avg_confidence
    
    def _ocr_paddleocr(self, image: np.ndarray) -> Tuple[str, float]:
        """Perform OCR using PaddleOCR"""
        results = self.paddleocr_reader.ocr(image, cls=True)
        
        if not results or not results[0]:
            return "", 0.0
        
        # Extract text and confidence
        texts = []
        confidences = []
        
        for line in results[0]:
            bbox, (text, conf) = line
            texts.append(text)
            confidences.append(conf)
        
        combined_text = ' '.join(texts)
        avg_confidence = sum(confidences) / len(confidences) if confidences else 0
        
        return combined_text, avg_confidence
    
    def _detect_language(self, text: str) -> str:
        """Detect language of extracted text"""
        if not LANGDETECT_AVAILABLE or not text.strip():
            return 'unknown'
        
        try:
            return langdetect.detect(text)
        except:
            return 'unknown'


class DocumentProcessor:
    """Advanced document processing for various formats"""
    
    def __init__(self):
        self.ocr_engine = MultiEngineOCR()
    
    def extract_from_pdf(self, file_path: str, method: str = 'auto') -> ExtractionResult:
        """Extract text from PDF using multiple methods"""
        start_time = datetime.now()
        
        if method == 'auto':
            # Try text extraction first, fall back to OCR if needed
            text_result = self._extract_pdf_text(file_path)
            if text_result and len(text_result.strip()) > 50:
                method = 'text'
            else:
                method = 'ocr'
        
        if method == 'text':
            text = self._extract_pdf_text(file_path)
            confidence = 0.95  # High confidence for text extraction
        elif method == 'ocr':
            text, confidence = self._extract_pdf_ocr(file_path)
        else:
            raise ValueError(f"Unknown PDF extraction method: {method}")
        
        processing_time = (datetime.now() - start_time).total_seconds()
        
        return ExtractionResult(
            text=text,
            confidence=confidence,
            language=self._detect_language(text),
            metadata={
                'method': method,
                'file_type': 'pdf',
                'file_size': os.path.getsize(file_path)
            },
            processing_time=processing_time,
            source_file=file_path,
            extraction_method='pdf',
            timestamp=datetime.now()
        )
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF using text-based methods"""
        texts = []
        
        # Try pdfplumber first (best for complex layouts)
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        texts.append(text)
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}")
        
        # Try PyMuPDF if available and pdfplumber didn't work well
        if not texts and PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(file_path)
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    text = page.get_text()
                    if text:
                        texts.append(text)
                doc.close()
            except Exception as e:
                logger.warning(f"PyMuPDF failed: {e}")
        
        # Fallback to PyPDF2
        if not texts:
            try:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            texts.append(text)
            except Exception as e:
                logger.warning(f"PyPDF2 failed: {e}")
        
        return '\n\n'.join(texts)
    
    def _extract_pdf_ocr(self, file_path: str) -> Tuple[str, float]:
        """Extract text from PDF using OCR"""
        try:
            if PYMUPDF_AVAILABLE:
                doc = fitz.open(file_path)
                texts = []
                confidences = []
                
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    # Convert page to image
                    mat = fitz.Matrix(2.0, 2.0)  # High resolution
                    pix = page.get_pixmap(matrix=mat)
                    img_data = pix.tobytes("png")
                    
                    # Convert to OpenCV format
                    nparr = np.frombuffer(img_data, np.uint8)
                    image = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                    
                    # Perform OCR
                    result = self.ocr_engine.extract_text(image)
                    texts.append(result.text)
                    confidences.append(result.confidence)
                
                doc.close()
                
                combined_text = '\n\n'.join(texts)
                avg_confidence = sum(confidences) / len(confidences) if confidences else 0
                
                return combined_text, avg_confidence
            else:
                raise Exception("PyMuPDF not available for PDF OCR")
                
        except Exception as e:
            logger.error(f"PDF OCR failed: {e}")
            return "", 0.0
    
    def extract_from_docx(self, file_path: str) -> ExtractionResult:
        """Extract text from Word document"""
        start_time = datetime.now()
        
        try:
            doc = Document(file_path)
            
            # Extract paragraphs
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Extract tables
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        table_texts.append(row_text)
            
            # Combine all text
            all_text = '\n'.join(paragraphs)
            if table_texts:
                all_text += '\n\nTables:\n' + '\n'.join(table_texts)
            
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return ExtractionResult(
                text=all_text,
                confidence=0.95,
                language=self._detect_language(all_text),
                metadata={
                    'paragraphs': len(paragraphs),
                    'tables': len(doc.tables),
                    'file_type': 'docx',
                    'file_size': os.path.getsize(file_path)
                },
                processing_time=processing_time,
                source_file=file_path,
                extraction_method='docx',
                timestamp=datetime.now()
            )
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise
    
    def extract_from_excel(self, file_path: str) -> ExtractionResult:
        """Extract text from Excel file"""
        start_time = datetime.now()
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            all_text = []
            total_cells = 0
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = [f"Sheet: {sheet_name}"]
                
                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                            total_cells += 1
                    
                    if row_values:
                        sheet_text.append(' | '.join(row_values))
                
                if len(sheet_text) > 1:  # More than just the header
                    all_text.extend(sheet_text)
                    all_text.append('')  # Empty line between sheets
            
            combined_text = '\n'.join(all_text)
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return ExtractionResult(
                text=combined_text,
                confidence=0.95,
                language=self._detect_language(combined_text),
                metadata={
                    'sheets': len(workbook.sheetnames),
                    'total_cells': total_cells,
                    'file_type': 'xlsx',
                    'file_size': os.path.getsize(file_path)
                },
                processing_time=processing_time,
                source_file=file_path,
                extraction_method='excel',
                timestamp=datetime.now()
            )
            
        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            raise
    
    def extract_from_powerpoint(self, file_path: str) -> ExtractionResult:
        """Extract text from PowerPoint presentation"""
        start_time = datetime.now()
        
        try:
            prs = Presentation(file_path)
            all_text = []
            slide_count = 0
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = [f"Slide {i}:"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the header
                    all_text.extend(slide_text)
                    all_text.append('')  # Empty line between slides
                    slide_count += 1
            
            combined_text = '\n'.join(all_text)
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return ExtractionResult(
                text=combined_text,
                confidence=0.95,
                language=self._detect_language(combined_text),
                metadata={
                    'slides': slide_count,
                    'file_type': 'pptx',
                    'file_size': os.path.getsize(file_path)
                },
                processing_time=processing_time,
                source_file=file_path,
                extraction_method='powerpoint',
                timestamp=datetime.now()
            )
            
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {e}")
            raise
    
    def _detect_language(self, text: str) -> str:
        """Detect language of text"""
        if not LANGDETECT_AVAILABLE or not text.strip():
            return 'unknown'
        
        try:
            return langdetect.detect(text)
        except:
            return 'unknown'


class AdvancedTextExtractor:
    """Main text extraction class coordinating all extraction methods"""
    
    def __init__(self, config_path: str = None):
        self.config = self._load_config(config_path)
        self.ocr_engine = MultiEngineOCR(self.config.get('ocr', {}))
        self.doc_processor = DocumentProcessor()
        self.db_path = self.config.get('database', {}).get('path', 'extraction_history.db')
        self._init_database()
        
        # Supported file types
        self.supported_image_types = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp'}
        self.supported_doc_types = {'.pdf', '.docx', '.xlsx', '.xls', '.pptx', '.txt'}
        
        logger.info("Advanced Text Extractor initialized")
    
    def _load_config(self, config_path: str) -> Dict[str, Any]:
        """Load configuration from file"""
        default_config = {
            'ocr': {
                'default_engine': 'auto',
                'preprocess': True,
                'confidence_threshold': 0.5
            },
            'database': {
                'path': 'extraction_history.db'
            },
            'batch': {
                'max_workers': 4,
                'chunk_size': 10
            }
        }
        
        if config_path and os.path.exists(config_path):
            try:
                with open(config_path, 'r') as f:
                    user_config = json.load(f)
                default_config.update(user_config)
            except Exception as e:
                logger.warning(f"Failed to load config: {e}")
        
        return default_config
    
    def _init_database(self):
        """Initialize extraction history database"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            cursor.execute('''
                CREATE TABLE IF NOT EXISTS extractions (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    file_path TEXT NOT NULL,
                    file_hash TEXT,
                    extraction_method TEXT,
                    text_content TEXT,
                    confidence REAL,
                    language TEXT,
                    metadata TEXT,
                    processing_time REAL,
                    timestamp TEXT,
                    UNIQUE(file_hash, extraction_method)
                )
            ''')
            
            conn.commit()
            conn.close()
            logger.info(f"Database initialized: {self.db_path}")
            
        except Exception as e:
            logger.error(f"Database initialization failed: {e}")
    
    def extract_from_file(self, file_path: str, method: str = 'auto', 
                         use_cache: bool = True) -> ExtractionResult:
        """Extract text from a file using the most appropriate method"""
        file_path = Path(file_path).resolve()
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Check cache first
        if use_cache:
            cached_result = self._get_cached_result(str(file_path))
            if cached_result:
                logger.info(f"Using cached result for {file_path.name}")
                return cached_result
        
        # Determine file type and extraction method
        file_ext = file_path.suffix.lower()
        
        if method == 'auto':
            if file_ext in self.supported_image_types:
                method = 'ocr'
            elif file_ext in self.supported_doc_types:
                method = 'document'
            else:
                # Try to detect file type by content
                try:
                    file_type = magic.from_file(str(file_path), mime=True)
                    if file_type.startswith('image/'):
                        method = 'ocr'
                    else:
                        method = 'document'
                except:
                    raise ValueError(f"Cannot determine extraction method for {file_ext}")
        
        # Perform extraction
        try:
            if method == 'ocr':
                result = self._extract_from_image(str(file_path))
            elif method == 'document':
                result = self._extract_from_document(str(file_path))
            else:
                raise ValueError(f"Unknown extraction method: {method}")
            
            # Cache the result
            if use_cache:
                self._cache_result(result)
            
            return result
            
        except Exception as e:
            logger.error(f"Extraction failed for {file_path}: {e}")
            raise
    
    def _extract_from_image(self, file_path: str) -> ExtractionResult:
        """Extract text from image file"""
        try:
            # Load image
            image = cv2.imread(file_path)
            if image is None:
                raise ValueError(f"Cannot load image: {file_path}")
            
            # Perform OCR
            result = self.ocr_engine.extract_text(image)
            result.source_file = file_path
            
            return result
            
        except Exception as e:
            logger.error(f"Image extraction failed: {e}")
            raise
    
    def _extract_from_document(self, file_path: str) -> ExtractionResult:
        """Extract text from document file"""
        file_ext = Path(file_path).suffix.lower()
        
        try:
            if file_ext == '.pdf':
                return self.doc_processor.extract_from_pdf(file_path)
            elif file_ext == '.docx':
                return self.doc_processor.extract_from_docx(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self.doc_processor.extract_from_excel(file_path)
            elif file_ext == '.pptx':
                return self.doc_processor.extract_from_powerpoint(file_path)
            elif file_ext == '.txt':
                return self._extract_from_text_file(file_path)
            else:
                raise ValueError(f"Unsupported document type: {file_ext}")
                
        except Exception as e:
            logger.error(f"Document extraction failed: {e}")
            raise
    
    def _extract_from_text_file(self, file_path: str) -> ExtractionResult:
        """Extract text from plain text file"""
        start_time = datetime.now()
        
        try:
            # Detect encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            # Read file
            with open(file_path, 'r', encoding=encoding) as f:
                text = f.read()
            
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return ExtractionResult(
                text=text,
                confidence=1.0,
                language=self._detect_language(text),
                metadata={
                    'encoding': encoding,
                    'file_type': 'txt',
                    'file_size': os.path.getsize(file_path)
                },
                processing_time=processing_time,
                source_file=file_path,
                extraction_method='text',
                timestamp=datetime.now()
            )
            
        except Exception as e:
            logger.error(f"Text file extraction failed: {e}")
            raise
    
    def batch_extract(self, file_paths: List[str], max_workers: int = None, 
                     progress_callback=None) -> List[ExtractionResult]:
        """Extract text from multiple files in parallel"""
        if max_workers is None:
            max_workers = self.config.get('batch', {}).get('max_workers', 4)
        
        results = []
        failed_files = []
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit all tasks
            future_to_file = {
                executor.submit(self.extract_from_file, file_path): file_path 
                for file_path in file_paths
            }
            
            # Process completed tasks
            for i, future in enumerate(as_completed(future_to_file)):
                file_path = future_to_file[future]
                
                try:
                    result = future.result()
                    results.append(result)
                    logger.info(f"Successfully extracted from {Path(file_path).name}")
                    
                except Exception as e:
                    failed_files.append((file_path, str(e)))
                    logger.error(f"Failed to extract from {file_path}: {e}")
                
                # Call progress callback
                if progress_callback:
                    progress_callback(i + 1, len(file_paths), file_path)
        
        if failed_files:
            logger.warning(f"Failed to process {len(failed_files)} files")
            for file_path, error in failed_files:
                logger.warning(f"  {file_path}: {error}")
        
        return results
    
    def _get_file_hash(self, file_path: str) -> str:
        """Calculate hash of file for caching"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def _cache_result(self, result: ExtractionResult):
        """Cache extraction result in database"""
        try:
            file_hash = self._get_file_hash(result.source_file)
            
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            cursor.execute('''
                INSERT OR REPLACE INTO extractions 
                (file_path, file_hash, extraction_method, text_content, confidence, 
                 language, metadata, processing_time, timestamp)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            ''', (
                result.source_file,
                file_hash,
                result.extraction_method,
                result.text,
                result.confidence,
                result.language,
                json.dumps(result.metadata),
                result.processing_time,
                result.timestamp.isoformat()
            ))
            
            conn.commit()
            conn.close()
            
        except Exception as e:
            logger.warning(f"Failed to cache result: {e}")
    
    def _get_cached_result(self, file_path: str) -> Optional[ExtractionResult]:
        """Get cached extraction result"""
        try:
            file_hash = self._get_file_hash(file_path)
            
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            cursor.execute('''
                SELECT * FROM extractions 
                WHERE file_hash = ? 
                ORDER BY timestamp DESC 
                LIMIT 1
            ''', (file_hash,))
            
            row = cursor.fetchone()
            conn.close()
            
            if row:
                return ExtractionResult(
                    text=row[3],
                    confidence=row[4],
                    language=row[5],
                    metadata=json.loads(row[6]),
                    processing_time=row[7],
                    source_file=file_path,
                    extraction_method=row[2],
                    timestamp=datetime.fromisoformat(row[8])
                )
            
        except Exception as e:
            logger.warning(f"Failed to get cached result: {e}")
        
        return None
    
    def _detect_language(self, text: str) -> str:
        """Detect language of text"""
        if not LANGDETECT_AVAILABLE or not text.strip():
            return 'unknown'
        
        try:
            return langdetect.detect(text)
        except:
            return 'unknown'
    
    def get_extraction_history(self, limit: int = 100) -> List[Dict[str, Any]]:
        """Get extraction history from database"""
        try:
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            cursor.execute('''
                SELECT * FROM extractions 
                ORDER BY timestamp DESC 
                LIMIT ?
            ''', (limit,))
            
            rows = cursor.fetchall()
            conn.close()
            
            history = []
            for row in rows:
                history.append({
                    'id': row[0],
                    'file_path': row[1],
                    'extraction_method': row[2],
                    'confidence': row[4],
                    'language': row[5],
                    'processing_time': row[7],
                    'timestamp': row[8],
                    'text_preview': row[3][:100] + '...' if len(row[3]) > 100 else row[3]
                })
            
            return history
            
        except Exception as e:
            logger.error(f"Failed to get extraction history: {e}")
            return []
    
    def export_results(self, results: List[ExtractionResult], 
                      output_path: str, format: str = 'json'):
        """Export extraction results to file"""
        try:
            if format.lower() == 'json':
                self._export_json(results, output_path)
            elif format.lower() == 'csv':
                self._export_csv(results, output_path)
            elif format.lower() == 'txt':
                self._export_txt(results, output_path)
            else:
                raise ValueError(f"Unsupported export format: {format}")
                
            logger.info(f"Results exported to {output_path}")
            
        except Exception as e:
            logger.error(f"Export failed: {e}")
            raise
    
    def _export_json(self, results: List[ExtractionResult], output_path: str):
        """Export results as JSON"""
        export_data = {
            'export_timestamp': datetime.now().isoformat(),
            'total_files': len(results),
            'results': []
        }
        
        for result in results:
            export_data['results'].append({
                'source_file': result.source_file,
                'text': result.text,
                'confidence': result.confidence,
                'language': result.language,
                'extraction_method': result.extraction_method,
                'processing_time': result.processing_time,
                'timestamp': result.timestamp.isoformat(),
                'metadata': result.metadata
            })
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(export_data, f, indent=2, ensure_ascii=False)
    
    def _export_csv(self, results: List[ExtractionResult], output_path: str):
        """Export results as CSV"""
        import csv
        
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            
            # Header
            writer.writerow([
                'Source File', 'Text', 'Confidence', 'Language', 
                'Extraction Method', 'Processing Time', 'Timestamp'
            ])
            
            # Data
            for result in results:
                writer.writerow([
                    result.source_file,
                    result.text.replace('\n', ' '),  # Single line for CSV
                    result.confidence,
                    result.language,
                    result.extraction_method,
                    result.processing_time,
                    result.timestamp.isoformat()
                ])
    
    def _export_txt(self, results: List[ExtractionResult], output_path: str):
        """Export results as plain text"""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"Text Extraction Results\n")
            f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"Total Files: {len(results)}\n")
            f.write("=" * 80 + "\n\n")
            
            for i, result in enumerate(results, 1):
                f.write(f"File {i}: {Path(result.source_file).name}\n")
                f.write(f"Method: {result.extraction_method}\n")
                f.write(f"Confidence: {result.confidence:.2f}\n")
                f.write(f"Language: {result.language}\n")
                f.write(f"Processing Time: {result.processing_time:.2f}s\n")
                f.write("-" * 40 + "\n")
                f.write(result.text)
                f.write("\n" + "=" * 80 + "\n\n")


if __name__ == "__main__":
    # Example usage
    extractor = AdvancedTextExtractor()
    
    # Extract from a single file
    # result = extractor.extract_from_file("example.pdf")
    # print(f"Extracted text: {result.text[:200]}...")
    # print(f"Confidence: {result.confidence:.2f}")
    # print(f"Language: {result.language}")
